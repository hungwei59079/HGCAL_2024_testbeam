import os

from pptx import Presentation
from pptx.util import Inches


# Function to add a slide with 6 images
def add_slide_with_images(prs, image_paths, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide.shapes.title
    title.text = title_text

    # Define positions for the 6 images (3x2 grid)
    positions = [
        (Inches(0.5), Inches(1.5), Inches(4), Inches(3)),  # Top-left
        (Inches(5), Inches(1.5), Inches(4), Inches(3)),  # Top-center
        (Inches(9.5), Inches(1.5), Inches(4), Inches(3)),  # Top-right
        (Inches(0.5), Inches(4.75), Inches(4), Inches(3)),  # Bottom-left
        (Inches(5), Inches(4.75), Inches(4), Inches(3)),  # Bottom-center
        (Inches(9.5), Inches(4.75), Inches(4), Inches(3)),  # Bottom-right
    ]

    for pos, image_path in zip(positions, image_paths):
        left, top, width, height = pos
        slide.shapes.add_picture(image_path, left, top, width, height)


# Create PowerPoint presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Set width for 16:9
prs.slide_height = Inches(7.5)  # Set height for 16:9
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Plot Overview"

os.chdir("./Trace_length_data/")

# Example image paths
image_sets = [
    [
        "Bias_-130_module_0_trace_correlation.png",
        "Bias_-130_module_1_trace_correlation.png",
        "Bias_-130_module_2_trace_correlation.png",
        "Bias_-130_module_3_trace_correlation.png",
        "Bias_-130_module_4_trace_correlation.png",
        "Bias_-130_module_5_trace_correlation.png",
    ],
    [
        "Bias_-300_module_0_trace_correlation.png",
        "Bias_-300_module_1_trace_correlation.png",
        "Bias_-300_module_2_trace_correlation.png",
        "Bias_-300_module_3_trace_correlation.png",
        "Bias_-300_module_4_trace_correlation.png",
        "Bias_-300_module_5_trace_correlation.png",
    ],
    [
        "Bias_-400_module_0_trace_correlation.png",
        "Bias_-400_module_1_trace_correlation.png",
        "Bias_-400_module_2_trace_correlation.png",
        "Bias_-400_module_3_trace_correlation.png",
        "Bias_-400_module_4_trace_correlation.png",
        "Bias_-400_module_5_trace_correlation.png",
    ],
    [
        "Bias_-500_module_0_trace_correlation.png",
        "Bias_-500_module_1_trace_correlation.png",
        "Bias_-500_module_2_trace_correlation.png",
        "Bias_-500_module_3_trace_correlation.png",
        "Bias_-500_module_4_trace_correlation.png",
        "Bias_-500_module_5_trace_correlation.png",
    ],
    [
        "Bias_-600_module_0_trace_correlation.png",
        "Bias_-600_module_1_trace_correlation.png",
        "Bias_-600_module_2_trace_correlation.png",
        "Bias_-600_module_3_trace_correlation.png",
        "Bias_-600_module_4_trace_correlation.png",
        "Bias_-600_module_5_trace_correlation.png",
    ],
    [
        "Bias_-700_module_0_trace_correlation.png",
        "Bias_-700_module_1_trace_correlation.png",
        "Bias_-700_module_2_trace_correlation.png",
        "Bias_-700_module_3_trace_correlation.png",
        "Bias_-700_module_4_trace_correlation.png",
        "Bias_-700_module_5_trace_correlation.png",
    ],
]

# Add slides for each set
for idx, image_paths in enumerate(image_sets):
    add_slide_with_images(prs, image_paths, f"Set {idx + 1}")

# Save the presentation
prs.save("plots_presentation.pptx")
print("Presentation saved as plots_presentation.pptx")
