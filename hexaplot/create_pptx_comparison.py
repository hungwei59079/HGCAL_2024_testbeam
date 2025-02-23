import os

from pptx import Presentation
from pptx.util import Inches


# Function to add a slide with 6 images
def add_slide_with_images(prs, image_paths, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide.shapes.title
    title.text = title_text

    # Define positions for the 6 images (3x2 grid)
    positions = [
        (Inches(0.5), Inches(1.5), Inches(4), Inches(3)),  # Top-left
        (Inches(5), Inches(1.5), Inches(4), Inches(3)),  # Top-center
        (Inches(9.5), Inches(1.5), Inches(4), Inches(3)),  # Top-right
        (Inches(0.5), Inches(4.75), Inches(4), Inches(3)),  # Bottom-left
        (Inches(5), Inches(4.75), Inches(4), Inches(3)),  # Bottom-center
        (Inches(9.5), Inches(4.75), Inches(4), Inches(3)),  # Bottom-right
    ]

    for pos, image_path in zip(positions, image_paths):
        left, top, width, height = pos
        slide.shapes.add_picture(image_path, left, top, width, height)


# Create PowerPoint presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Set width for 16:9
prs.slide_height = Inches(7.5)  # Set height for 16:9
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Plot Overview"

os.chdir("../root_files/histograms/")
# Example image paths
image_sets = [
    [
        "mod_0_ch_74-130.png",
        "mod_0_ch_74-300.png",
        "mod_0_ch_74-400.png",
        "mod_0_ch_74-500.png",
        "mod_0_ch_74-600.png",
        "mod_0_ch_74-700.png",
    ],
    [
        "mod_5_ch_9-130.png",
        "mod_5_ch_9-300.png",
        "mod_5_ch_9-400.png",
        "mod_5_ch_9-500.png",
        "mod_5_ch_9-600.png",
        "mod_5_ch_9-700.png",
    ],
]

# Add slides for each set
for idx, image_paths in enumerate(image_sets):
    add_slide_with_images(prs, image_paths, f"Set {idx + 1}")

# Save the presentation
prs.save("Comparison.pptx")
print("Presentation saved as plots_presentation.pptx")
